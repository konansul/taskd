import re
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from googletrans import Translator
from datetime import date


from .chart import add_chart
from .prompt import generate_image_hf





def add_title_slide(prs, slide):
    today = date.today()
    formatted_date = today.strftime("%d/%m/%Y")
    s = prs.slides[0]

    title_shape = s.shapes.title
    if title_shape and title_shape.text_frame:
        p = title_shape.text_frame.paragraphs[0]
        if p.runs:
            p.runs[0].text = slide['title']
        else:
            p.text = slide['title']

    for placeholder in s.placeholders:
        if "Tarix" in placeholder.text:
            text_frame = placeholder.text_frame
            p = text_frame.paragraphs[0]
            if p.runs:
                p.runs[0].text = f"Tarix: {formatted_date}"
            else:
                p.text = f"Tarix: {formatted_date}"
            break


def add_intro_slide(prs, slide):
    s = prs.slides[1]  # Target the second slide
    vertical_margin = Inches(0.2)

    for shape in s.shapes:
        if shape.has_text_frame:
            text = shape.text.strip()
            if text == "Layihənin məzmunu":

                if 'summary' in slide and slide['summary']:
                    new_text_box_width = shape.width
                    left = shape.left
                    top = shape.top + shape.height + vertical_margin

                    initial_height_for_autosize = Inches(3)

                    txBox = s.shapes.add_textbox(left, top, new_text_box_width, initial_height_for_autosize)
                    tf = txBox.text_frame

                    tf.word_wrap = True
                    tf.margin_left = 0
                    tf.margin_right = 0
                    tf.margin_top = 0
                    tf.margin_bottom = 0

                    p_content = tf.paragraphs[0]
                    p_content.clear()
                    run_content = p_content.add_run()
                    run_content.text = slide['summary']
                    run_content.font.size = Pt(17)
                    run_content.font.bold = False

                else:
                    print("Summary data is missing or empty.")

            elif text == "Məqsəd":

                if 'aim' in slide and slide['aim']:
                    new_text_box_width = shape.width
                    left = shape.left
                    top = shape.top + shape.height + vertical_margin

                    initial_height_for_autosize = Inches(3)

                    txBox = s.shapes.add_textbox(left, top, new_text_box_width, initial_height_for_autosize)
                    tf = txBox.text_frame

                    tf.word_wrap = True
                    tf.margin_left = 0
                    tf.margin_right = 0
                    tf.margin_top = 0
                    tf.margin_bottom = 0

                    p_content = tf.paragraphs[0]
                    p_content.clear()
                    run_content = p_content.add_run()
                    run_content.text = slide['aim']
                    run_content.font.size = Pt(17)
                    run_content.font.bold = False



                else:
                    print("Aim data is missing or empty.")
        else:
            print("    Does not have a text frame.")


def safe_float_conversion(values):
    result = []
    for v in values:
        try:
            if isinstance(v, str):
                v = v.strip()
                # Handle percentages
                if v.endswith('%'):
                    result.append(float(v[:-1]))
                    continue
                # Handle "Vmax" or other common non-numeric placeholders by treating as 0 or skipping
                # For now, 0.0 is safer to keep alignment with X axis
            result.append(float(v))
        except (ValueError, TypeError):
            print(f"Warning: Could not convert '{v}' to float. Using 0.0 instead.")
            result.append(0.0)
    return result


def add_main_slide(prs, slide):
    layout = prs.slide_layouts[12]
    s = prs.slides.add_slide(layout)

    s.shapes.title.text = slide['title']
    current_point_idx = 1

    for shape in s.shapes:
        if not shape.is_placeholder or shape == s.shapes.title:
            continue

        current_text_in_shape = shape.text.strip().lower()

        if current_point_idx <= 4 and (current_text_in_shape in ("", "xxxx") or True):
            point = slide.get(f'point{current_point_idx}', '')

            tf = shape.text_frame
            tf.clear()

            if point:
                p = tf.paragraphs[0]
                p.text = point
                p.font.size = Pt(17)
            else:
                pass

            current_point_idx += 1

    visual = slide.get('visual', {})
    if visual.get('type') and visual['type'] != 'none':
        visual_slide_layout = prs.slide_layouts[3]
        visual_s = prs.slides.add_slide(visual_slide_layout)
        visual_s.shapes.title.text = f"{slide['title']} - {visual.get('title', 'Visual')}"

        if visual["type"] in ["bar", "line"]:
            # Use safe conversion for Y values
            y_values = safe_float_conversion(visual["y"])
            add_chart(visual_s, visual["type"], visual["title"],
                      visual["x"], y_values,
                      visual.get("xlabel", ""), visual.get("ylabel", ""))
        elif visual["type"] == "pie":
            sizes = safe_float_conversion(visual["sizes"])
            add_chart(visual_s, visual["type"], visual["title"],
                      labels=visual['labels'], sizes=sizes)
        elif visual["type"] == "image":
            try:
                image_path = f"generated_image_{visual['title'][:10].replace(' ', '_')}.png"
                translator = Translator()
                translated = translator.translate(visual["description"], src='az', dest='en')
                english_description = translated.text
                generate_image_hf(english_description, image_path)

                try:
                    placeholder = visual_s.placeholders[1]
                    left, top, width, height = placeholder.left, placeholder.top, placeholder.width, placeholder.height
                except IndexError:
                    left, top, width, height = Inches(1), Inches(2), Inches(8), Inches(4)

                visual_s.shapes.add_picture(image_path, left, top, width=width, height=height)
            except Exception as e:
                insert_text_or_fallback(visual_s, f"[Şəkil təsviri: {visual.get('description','')}]")
        else:
            insert_text_or_fallback(visual_s, f"[Vizual növü '{visual['type']}' hələ dəstəklənmir]")

def insert_text_or_fallback(slide, text):
    try:
        placeholder = slide.placeholders[1]  # second placeholder
        placeholder.text = text
    except IndexError:
        textbox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(6), Inches(1))
        textbox.text_frame.text = text


def add_recommendation_slide(prs, slide):
    layout = prs.slide_layouts[3]
    s = prs.slides.add_slide(layout)
    title_shape = s.shapes.title

    if title_shape:
        title_shape.text = "Növbəti addımlar"

    target_tf = None

    for shape in s.shapes:

        if not shape.has_text_frame:
            continue

        paragraphs = shape.text_frame.paragraphs
        if not paragraphs:
            continue

        first_text = paragraphs[0].text.strip()

        if first_text == "":
            target_tf = shape.text_frame

    if target_tf:

        target_tf.clear()

        for i in range(1, 6):
            rec = slide.get(f'recommendation{i}', None)
            if rec:
                p = target_tf.add_paragraph()
                p.text = rec
                p.level = 1
                p.font.color.rgb = RGBColor(0, 0, 0)
                p.font.size = Pt(21)

                print(f"Added recommendation{i}: {rec}")
            else:
                print(f"No data for recommendation{i}.")


def delete_slide(prs, slide_index):
    slide_id = prs.slides._sldIdLst[slide_index]
    prs.slides._sldIdLst.remove(slide_id)


def generate_pptx(slides, output_filename="presentation.pptx"):
    # Get the directory where this script is located
    import os
    template_dir = os.path.dirname(os.path.abspath(__file__))
    template_path = os.path.join(template_dir, "..", "format_new.pptx")
    prs = Presentation(template_path)

    for slide in slides:
        t = slide.get('type')
        if t == 'title':
            add_title_slide(prs, slide)
        elif t == 'intro':
            add_intro_slide(prs, slide)
        elif t == 'main':
            add_main_slide(prs, slide)
        elif t == 'recommendation':
            add_recommendation_slide(prs, slide)

    delete_slide(prs, 3)
    delete_slide(prs, 2)

    prs.save(output_filename)
    print(f"Presentation saved as '{output_filename}'")


def parse_gpt_response(response_text):
    try:
        # Check if response is an error message
        if not response_text or not isinstance(response_text, str):
            raise ValueError("Empty or invalid response received from presentation generator.")
        
        text = response_text.strip()
        
        # Check for error messages from the API
        if text.startswith("Error:") or text.startswith("Content generation failed"):
            raise ValueError(f"API error: {text}")
        
        # Check if response is too short (likely an error or incomplete)
        if len(text) < 50:
            raise ValueError(f"Response too short or incomplete: {text[:100]}")
        
        # Remove common markdown code fences (```json ... ```)
        # Handle multiple formats: ```json ... ```, ``` ... ```, etc.
        if "```" in text:
            # Remove opening ```json or ```
            text = re.sub(r"^```(?:json|JSON)?\s*", "", text, flags=re.MULTILINE)
            # Remove closing ```
            text = re.sub(r"```\s*$", "", text, flags=re.MULTILINE)
            text = text.strip()

        # Try to find JSON array first (most common case)
        # Look for complete JSON array: [...]
        json_array_pattern = re.compile(r'\[[\s\S]*\]', re.DOTALL)
        json_object_pattern = re.compile(r'\{[\s\S]*\}', re.DOTALL)
        
        cleaned_text = ""
        
        # First try to find a JSON array
        array_match = json_array_pattern.search(text)
        if array_match:
            cleaned_text = array_match.group(0).strip()
        else:
            # If no array found, try to find a JSON object
            object_match = json_object_pattern.search(text)
            if object_match:
                cleaned_text = object_match.group(0).strip()
            else:
                # If no JSON structure found, use the whole text
                cleaned_text = text

        if not cleaned_text:
            raise ValueError(f"Empty response after cleaning. Original: {text[:200]}...")
        
        # Try to parse JSON
        try:
            slides = json.loads(cleaned_text)
        except json.JSONDecodeError as json_err:
            # If parsing fails, try to fix common issues
            # Remove trailing commas
            cleaned_text = re.sub(r',\s*}', '}', cleaned_text)
            cleaned_text = re.sub(r',\s*]', ']', cleaned_text)
            
            # Try to extract just the array part more carefully
            # Find the first [ and last ]
            first_bracket = cleaned_text.find('[')
            last_bracket = cleaned_text.rfind(']')
            
            if first_bracket != -1 and last_bracket != -1 and last_bracket > first_bracket:
                cleaned_text = cleaned_text[first_bracket:last_bracket + 1]
            
            try:
                slides = json.loads(cleaned_text)
            except json.JSONDecodeError as json_err2:
                # Last resort: show more details about the error
                error_pos = json_err2.pos if hasattr(json_err2, 'pos') else 'unknown'
                error_msg = json_err2.msg if hasattr(json_err2, 'msg') else str(json_err2)
                preview = cleaned_text[max(0, int(error_pos) - 50):int(error_pos) + 50] if isinstance(error_pos, int) else cleaned_text[:200]
                raise ValueError(
                    f"Invalid JSON format at position {error_pos}: {error_msg}\n"
                    f"Preview: ...{preview}...\n"
                    f"Full response length: {len(response_text)} characters"
                )

    except json.JSONDecodeError as e:
        # This should not happen now, but keep it as fallback
        error_pos = e.pos if hasattr(e, 'pos') else 'unknown'
        error_msg = e.msg if hasattr(e, 'msg') else str(e)
        raise ValueError(f"Invalid JSON format at position {error_pos}: {error_msg}")

    # Basic validation: check it's a list, and each slide has required fields depending on type
    if not isinstance(slides, list):
        raise ValueError("Expected a JSON array (list) of slides.")

    for i, slide in enumerate(slides):
        if not isinstance(slide, dict):
            raise ValueError(f"Slide {i} is not a JSON object.")
        slide_type = slide.get("type")
        if slide_type not in {"title", "intro", "main", "recommendation"}:
            raise ValueError(f"Slide {i} has unknown type: {slide_type}")

        # Example minimal checks per slide type
        if slide_type == "title":
            if "title" not in slide:
                raise ValueError(f"Slide {i} of type 'title' missing required fields.")
        elif slide_type == "intro":
            if "aim" not in slide or "summary" not in slide:
                raise ValueError(f"Slide {i} of type 'intro' missing required fields.")
        elif slide_type == "main":
            required_fields = {"title", "point1", "point2", "point3", "point4", "visual"}
            if not required_fields.issubset(slide.keys()):
                missing = required_fields - set(slide.keys())
                raise ValueError(f"Slide {i} of type 'main' missing fields: {missing}")
            # Visual should be a dict
            if not isinstance(slide["visual"], dict):
                raise ValueError(f"Slide {i} 'visual' must be an object.")
        elif slide_type == "recommendation":
            # At least 4 recommendations required
            rec_fields = [f"recommendation{i}" for i in range(1, 6)]
            present_recs = [f for f in rec_fields if f in slide]
            if len(present_recs) < 4:
                raise ValueError(f"Slide {i} of type 'recommendation' requires at least 4 recommendations.")

    return slides